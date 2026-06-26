import os
import csv
import json
import logging
import re
from typing import List, Dict, Any

# Third-party libraries (imported safely or assumed available)
import pypdf
try:
    import docx
except ImportError:
    docx = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

from backend.ocr.service import ocr_service
from backend.core.config import settings

logger = logging.getLogger(__name__)

# Stopwords for keyword extraction
STOPWORDS = {
    "the", "a", "an", "and", "or", "but", "if", "then", "else", "when", "at", "by", 
    "from", "for", "in", "out", "on", "off", "over", "under", "again", "further", 
    "then", "once", "here", "there", "when", "where", "why", "how", "all", "any", 
    "both", "each", "few", "more", "most", "other", "some", "such", "no", "nor", 
    "not", "only", "own", "same", "so", "than", "too", "very", "s", "t", "can", 
    "will", "just", "don", "should", "now", "of", "to", "is", "was", "were", "are", 
    "has", "have", "had", "this", "that", "these", "those", "with", "about", "into"
}

class DocumentParser:
    def parse(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Parses a document based on its file extension.
        
        Returns a dict: 
        {
            "pages": [{"text": str, "page_number": int, "metadata": dict}],
            "metadata": {
                "author": str,
                "creation_date": str,
                "modified_date": str,
                "page_count": int,
                "language": str,
                "category": str,
                "keywords": List[str],
                "topics": List[str]
            }
        }
        """
        file_type = file_type.lower()
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        logger.info(f"Parsing file {file_path} of type {file_type}")
        
        pages = []
        author = "Unknown"
        creation_date = "Unknown"
        modified_date = "Unknown"

        if file_type == "pdf":
            pages, author, creation_date, modified_date = self._parse_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            pages, author, creation_date, modified_date = self._parse_docx(file_path)
        elif file_type == "xlsx":
            pages, author, creation_date, modified_date = self._parse_xlsx(file_path)
        elif file_type == "csv":
            pages = self._parse_csv(file_path)
        elif file_type in ["pptx", "ppt"]:
            pages, author, creation_date, modified_date = self._parse_pptx(file_path)
        elif file_type in ["txt", "md", "markdown", "json", "html"]:
            pages = self._parse_text_file(file_path, file_type)
        elif file_type in ["png", "jpg", "jpeg", "webp", "tiff", "bmp"]:
            pages = self._parse_image_file(file_path)
        else:
            pages = self._parse_text_file(file_path, "txt")

        # Extract OS dates as fallback
        if creation_date == "Unknown" or modified_date == "Unknown":
            try:
                stat = os.stat(file_path)
                if creation_date == "Unknown":
                    import datetime
                    creation_date = datetime.datetime.fromtimestamp(stat.st_ctime).isoformat()
                if modified_date == "Unknown":
                    import datetime
                    modified_date = datetime.datetime.fromtimestamp(stat.st_mtime).isoformat()
            except Exception:
                pass

        # Compile full text to analyze language and topics
        full_text = "\n".join([p["text"] for p in pages])
        
        # 1. Detect Language (Telugu / Hindi / English)
        language = self._detect_language(full_text)
        
        # 2. Extract Keywords and Topics
        keywords, topics = self._extract_keywords_and_topics(full_text)

        # 3. Classify Category
        category = self._classify_category(full_text, file_type)

        doc_metadata = {
            "author": author,
            "creation_date": creation_date,
            "modified_date": modified_date,
            "page_count": len(pages),
            "language": language,
            "category": category,
            "keywords": keywords,
            "topics": topics
        }

        return {
            "pages": pages,
            "metadata": doc_metadata
        }

    def _detect_language(self, text: str) -> str:
        """Determines document language using unicode block checks (offline)."""
        if not text:
            return "English"
            
        # Check Devanagari range (Hindi)
        hindi_chars = len(re.findall(r"[\u0900-\u097F]", text))
        # Check Telugu range
        telugu_chars = len(re.findall(r"[\u0C00-\u0C7F]", text))
        
        total = len(text)
        if total == 0:
            return "English"
            
        # Heuristics
        if telugu_chars > 20 or (telugu_chars / total) > 0.05:
            return "Telugu"
        elif hindi_chars > 20 or (hindi_chars / total) > 0.05:
            return "Hindi"
            
        return "English"

    def _classify_category(self, text: str, file_type: str) -> str:
        """Heuristics to categorize documents."""
        text_lower = text.lower()
        
        if "invoice" in text_lower or "receipt" in text_lower or "billing" in text_lower or "total amount" in text_lower:
            return "Invoice/Receipt"
        elif "abstract" in text_lower or "introduction" in text_lower or "conclusion" in text_lower or "references" in text_lower:
            return "Research Paper"
        elif "manual" in text_lower or "user guide" in text_lower or "installation" in text_lower or "step" in text_lower:
            return "User Manual"
        elif "policy" in text_lower or "corporate" in text_lower or "conduct" in text_lower or "hr agreement" in text_lower:
            return "Corporate Policy"
        elif "api" in text_lower or "specification" in text_lower or "developer" in text_lower or "function" in text_lower:
            return "Technical Documentation"
        
        # File type fallbacks
        if file_type in ["xlsx", "csv"]:
            return "Financial Spreadsheet"
        elif file_type in ["pptx", "ppt"]:
            return "Presentation Deck"
            
        return "General Document"

    def _extract_keywords_and_topics(self, text: str) -> tuple[List[str], List[str]]:
        """Extracts significant keywords and general topics from text content."""
        if not text:
            return [], []
            
        # Clean text
        words = re.findall(r"\b[a-zA-Z]{4,15}\b", text.lower())
        filtered_words = [w for w in words if w not in STOPWORDS]
        
        # Count frequencies
        freq: Dict[str, int] = {}
        for w in filtered_words:
            freq[w] = freq.get(w, 0) + 1
            
        # Sort and select top keywords
        sorted_words = sorted(freq.keys(), key=lambda x: freq[x], reverse=True)
        keywords = sorted_words[:6]
        
        # Build topics based on keyword categories
        topics = []
        if any(w in sorted_words[:20] for w in ["finance", "cost", "revenue", "budget", "tax", "billing"]):
            topics.append("Financial Management")
        if any(w in sorted_words[:20] for w in ["algorithm", "data", "model", "python", "software", "api"]):
            topics.append("Software Engineering")
        if any(w in sorted_words[:20] for w in ["hiring", "policy", "leave", "employee", "salary", "recruitment"]):
            topics.append("Human Resources")
        if any(w in sorted_words[:20] for w in ["patent", "research", "results", "analysis", "experiments"]):
            topics.append("Scientific Research")

        if not topics:
            # Fallback topics from top keywords
            topics = [w.capitalize() for w in keywords[:2]]
            
        return keywords, topics

    def _parse_pdf(self, file_path: str) -> tuple[List[Dict[str, Any]], str, str, str]:
        pages = []
        author, creation_date, modified_date = "Unknown", "Unknown", "Unknown"
        
        try:
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                num_pages = len(reader.pages)
                
                # Extract PDF properties
                meta = reader.metadata
                if meta:
                    author = meta.get("/Author", "Unknown")
                    creation_date = meta.get("/CreationDate", "Unknown")
                    modified_date = meta.get("/ModDate", "Unknown")
                
                for idx in range(num_pages):
                    page = reader.pages[idx]
                    page_num = idx + 1
                    text = page.extract_text() or ""
                    
                    # Check if the page is scanned
                    if ocr_service.is_text_scanned(text):
                        logger.info(f"Page {page_num} in PDF is scanned. Running local OCR...")
                        if fitz:
                            ocr_text = self._ocr_pdf_page(file_path, idx)
                            if ocr_text:
                                text = ocr_text
                        else:
                            logger.warning("PyMuPDF (fitz) is not installed. Scanned PDF page OCR skipped.")
                            
                    pages.append({
                        "text": text.strip(),
                        "page_number": page_num,
                        "metadata": {"source_type": "pdf"}
                    })
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            raise Exception(f"PDF parsing error: {e}") from e
            
        return pages, str(author), str(creation_date), str(modified_date)

    def _ocr_pdf_page(self, pdf_path: str, page_index: int) -> str:
        if not fitz:
            return ""
            
        temp_img_path = ""
        try:
            doc = fitz.open(pdf_path)
            page = doc[page_index]
            matrix = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=matrix)
            
            temp_img_name = f"temp_ocr_{os.path.basename(pdf_path)}_{page_index}.png"
            temp_img_path = os.path.join(settings.TEMP_DIR, temp_img_name)
            pix.save(temp_img_path)
            
            ocr_text = ocr_service.extract_text(temp_img_path)
            return ocr_text
        except Exception as e:
            logger.error(f"Failed to OCR PDF page {page_index}: {e}")
            return ""
        finally:
            if temp_img_path and os.path.exists(temp_img_path):
                try:
                    os.remove(temp_img_path)
                except Exception:
                    pass

    def _parse_docx(self, file_path: str) -> tuple[List[Dict[str, Any]], str, str, str]:
        if not docx:
            raise ImportError("python-docx package is not installed.")
            
        doc = docx.Document(file_path)
        author, creation_date, modified_date = "Unknown", "Unknown", "Unknown"
        
        # Read core properties
        try:
            props = doc.core_properties
            if props.author:
                author = props.author
            if props.created:
                creation_date = props.created.isoformat()
            if props.modified:
                modified_date = props.modified.isoformat()
        except Exception:
            pass

        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
                
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    full_text.append(" | ".join(row_text))

        text_content = "\n".join(full_text)
        pages = [{
            "text": text_content.strip(),
            "page_number": 1,
            "metadata": {"source_type": "docx"}
        }]
        return pages, author, creation_date, modified_date

    def _parse_xlsx(self, file_path: str) -> tuple[List[Dict[str, Any]], str, str, str]:
        if not openpyxl:
            raise ImportError("openpyxl package is not installed.")
            
        wb = openpyxl.load_workbook(file_path, data_only=True)
        author, creation_date, modified_date = "Unknown", "Unknown", "Unknown"
        
        # Read workbook properties
        try:
            props = wb.properties
            if props.creator:
                author = props.creator
            if props.created:
                creation_date = props.created.isoformat()
            if props.modified:
                modified_date = props.modified.isoformat()
        except Exception:
            pass

        sheets_data = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            
            for row in sheet.iter_rows(values_only=True):
                row_str = [str(cell).strip() if cell is not None else "" for cell in row]
                if any(row_str):
                    sheet_rows.append(" | ".join(row_str))
                    
            if sheet_rows:
                sheets_data.append({
                    "text": f"Sheet: {sheet_name}\n" + "\n".join(sheet_rows),
                    "page_number": len(sheets_data) + 1,
                    "metadata": {"source_type": "excel", "sheet_name": sheet_name}
                })
                
        return sheets_data, author, creation_date, modified_date

    def _parse_csv(self, file_path: str) -> List[Dict[str, Any]]:
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                if row:
                    rows.append(" | ".join(row))
                    
        return [{
            "text": "\n".join(rows),
            "page_number": 1,
            "metadata": {"source_type": "csv"}
        }]

    def _parse_pptx(self, file_path: str) -> tuple[List[Dict[str, Any]], str, str, str]:
        if not Presentation:
            raise ImportError("python-pptx package is not installed.")
            
        prs = Presentation(file_path)
        author, creation_date, modified_date = "Unknown", "Unknown", "Unknown"
        
        try:
            props = prs.core_properties
            if props.author:
                author = props.author
            if props.created:
                creation_date = props.created.isoformat()
            if props.modified:
                modified_date = props.modified.isoformat()
        except Exception:
            pass

        slides = []
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                    
            slides.append({
                "text": "\n".join(slide_text),
                "page_number": idx + 1,
                "metadata": {"source_type": "powerpoint"}
            })
            
        return slides, author, creation_date, modified_date

    def _parse_text_file(self, file_path: str, file_type: str) -> List[Dict[str, Any]]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
            
        if file_type == "json":
            try:
                parsed_json = json.loads(text)
                text = json.dumps(parsed_json, indent=2)
            except Exception:
                pass
                
        return [{
            "text": text.strip(),
            "page_number": 1,
            "metadata": {"source_type": file_type}
        }]

    def _parse_image_file(self, file_path: str) -> List[Dict[str, Any]]:
        text = ocr_service.extract_text(file_path)
        return [{
            "text": text,
            "page_number": 1,
            "metadata": {"source_type": "image", "is_image_ocr": True}
        }]

document_parser = DocumentParser()
